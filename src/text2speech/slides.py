"""Extract text from presentation slides (PPTX, PDF)."""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class Slide:
    index: int          # 1-based
    title: str
    body: str           # all non-title text joined
    notes: str          # presenter notes
    raw_lines: list[str] = field(default_factory=list)
    images: list[bytes] = field(default_factory=list)   # raw image blobs (PNG/JPEG)

    def spoken_text(self, include_notes: bool = False) -> str:
        """Return the text that should be spoken for this slide."""
        parts: list[str] = []
        if self.title:
            parts.append(self.title)
        if self.body:
            parts.append(self.body)
        if include_notes and self.notes:
            parts.append(self.notes)
        return "\n".join(parts).strip()

    def display_title(self) -> str:
        return self.title or f"Slide {self.index}"


def _clean(text: str) -> str:
    return " ".join(text.split()).strip()


# ── PPTX ─────────────────────────────────────────────────────────────────────

def read_pptx(path: Path) -> list[Slide]:
    try:
        from pptx import Presentation  # type: ignore[import]
        from pptx.util import Pt       # type: ignore[import]
    except ImportError:
        raise RuntimeError(
            "python-pptx is required for .pptx files.\n"
            "Install: uv pip install --default-index https://pypi.org/simple python-pptx"
        )

    prs = Presentation(str(path))
    slides: list[Slide] = []

    MSO_PICTURE = 13  # MSO_SHAPE_TYPE.PICTURE

    for i, slide in enumerate(prs.slides, 1):
        title = ""
        body_lines: list[str] = []
        notes_text = ""
        image_blobs: list[bytes] = []

        for shape in slide.shapes:
            # Collect picture blobs
            if shape.shape_type == MSO_PICTURE:
                try:
                    image_blobs.append(shape.image.blob)
                except Exception:
                    pass
                continue

            # Named chrome shapes — skip them so they don't pollute bullet content
            shape_name = getattr(shape, "name", "")
            if shape_name in ("slide_counter", "slide_tag"):
                continue
            if shape_name == "slide_title":
                if shape.has_text_frame and not title:
                    title = _clean(shape.text_frame.text)
                continue

            # Some shapes embed an image via a picture placeholder (ph_type 18)
            try:
                ph = shape.placeholder_format
            except (ValueError, AttributeError):
                ph = None
            ph_type = getattr(ph, "type", None)

            if ph_type == 18:   # OBJECT / picture placeholder
                try:
                    image_blobs.append(shape.image.blob)
                except Exception:
                    pass
                continue

            if not shape.has_text_frame:
                continue

            shape_text = "\n".join(
                _clean(para.text)
                for para in shape.text_frame.paragraphs
                if para.text.strip()
            )
            if not shape_text:
                continue

            if ph_type in (1, 13, 15):  # TITLE, CENTER_TITLE, SUBTITLE
                if not title:
                    title = _clean(shape.text_frame.text)
                    continue
            body_lines.append(shape_text)

        # Presenter notes
        if slide.has_notes_slide:
            nf = slide.notes_slide.notes_text_frame
            notes_text = _clean(nf.text) if nf else ""

        slides.append(Slide(
            index=i,
            title=title,
            body="\n".join(body_lines),
            notes=notes_text,
            raw_lines=[title] + body_lines,
            images=image_blobs,
        ))

    return slides


# ── PDF ──────────────────────────────────────────────────────────────────────

def read_pdf(path: Path) -> list[Slide]:
    try:
        import pypdf  # type: ignore[import]
    except ImportError:
        raise RuntimeError(
            "pypdf is required for .pdf files.\n"
            "Install: uv pip install --default-index https://pypi.org/simple pypdf"
        )

    reader = pypdf.PdfReader(str(path))
    slides: list[Slide] = []

    for i, page in enumerate(reader.pages, 1):
        raw = page.extract_text() or ""
        lines = [_clean(ln) for ln in raw.splitlines() if ln.strip()]
        title = lines[0] if lines else ""
        body = "\n".join(lines[1:]) if len(lines) > 1 else ""
        slides.append(Slide(
            index=i,
            title=title,
            body=body,
            notes="",
            raw_lines=lines,
        ))

    return slides


# ── dispatcher ────────────────────────────────────────────────────────────────

def read_slides(path: Path) -> list[Slide]:
    """Load slides from a .pptx or .pdf file."""
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return read_pptx(path)
    if suffix in (".pdf",):
        return read_pdf(path)
    raise ValueError(
        f"Unsupported format: {suffix!r}. Supported: .pptx, .pdf"
    )
