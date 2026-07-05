"""Build a styled 16:9 PPTX from a JSON slide plan.

The public API is :func:`build_pptx`.  All other names are private helpers.
"""

from __future__ import annotations
import io
import warnings
from pathlib import Path

# ── Slide dimensions (16:9 widescreen, EMU units) ─────────────────────────────
from pptx.util import Inches, Pt, Emu
W = Inches(13.333)
H = Inches(7.5)

# ── Colour palette — dark charcoal + forest green ──────────────────────────────
from pptx.dml.color import RGBColor

C_BG       = RGBColor(0xF6, 0xFA, 0xF7)  # near-white, faint green tint
C_HEADER   = RGBColor(0x14, 0x1A, 0x14)  # dark charcoal-forest
C_HDR_MID  = RGBColor(0x22, 0x2E, 0x22)  # lighter charcoal (depth band at top)
C_SKY      = RGBColor(0x16, 0xA3, 0x4A)  # emerald green — single accent colour
C_SKY_SOFT = RGBColor(0xDC, 0xFC, 0xE7)  # very light green (image panel fill)
C_TITLE    = RGBColor(0xFF, 0xFF, 0xFF)  # white
C_TEXT     = RGBColor(0x14, 0x23, 0x17)  # dark forest body text
C_SUB      = RGBColor(0x4A, 0x64, 0x4E)  # muted forest-green sub-bullet
C_DIM      = RGBColor(0x86, 0xAA, 0x8A)  # muted counter / decorative
C_RULE     = RGBColor(0xD1, 0xE8, 0xD4)  # light green rule / progress track
C_IMGBG    = RGBColor(0xFF, 0xFF, 0xFF)  # white image panel interior
C_IMGBDR   = RGBColor(0x16, 0xA3, 0x4A)  # emerald green image border

# backward-compat aliases used in canvas_video.py
C_ACCENT  = C_SKY
C_ACCENT2 = RGBColor(0x86, 0xEF, 0xAC)  # light green

# ── Layout constants (Inches) ──────────────────────────────────────────────────
HDR_H    = Inches(1.92)     # header height (tighter = more content space)
ACCENT_W = Inches(0.06)     # slim left accent bar
MARGIN_X = Inches(0.62)     # left margin
MARGIN_Y = Inches(0.26)     # top margin inside header
GUTTER   = Inches(0.22)     # gap between accent bar and bullets
IMG_X    = Inches(8.80)     # image panel left edge
IMG_W    = Inches(4.20)     # image panel width
IMG_TOP  = Inches(2.08)     # image panel top
IMG_BTM  = Inches(7.15)     # image panel bottom


# ── Helpers ────────────────────────────────────────────────────────────────────

def _solid(shape, color: RGBColor) -> None:
    """Fill a shape with a solid colour and remove its border."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _textbox(slide, left, top, width, height):
    """Add a transparent-background text box."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.fill.background()
    tb.line.fill.background()
    return tb


def _para(tf, text: str, size_pt: int, bold: bool = False,
          color: RGBColor = C_TEXT, space_before_pt: int = 0,
          align=None) -> None:
    """Append a paragraph to text-frame *tf*."""
    from pptx.enum.text import PP_ALIGN
    p = tf.add_paragraph()
    p.text = text
    p.space_before = Pt(space_before_pt)
    if align:
        p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _set_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = RECTANGLE
    _solid(shape, color)
    return shape


def _extract_pdf_page_image(pdf_path: Path, page_num: int) -> bytes | None:
    """Return PNG bytes for the best figure on `page_num` of `pdf_path`.

    Strategy:
    1. Extract the largest embedded raster image from the page (fast, exact).
    2. If none found, render the whole page via pypdfium2 (handles vector graphics).
    3. Returns None if both approaches fail.
    """
    try:
        import pypdf
        from PIL import Image as PILImage
    except ImportError:
        return None

    warnings.filterwarnings("ignore", message=".*Lookup Table.*")
    reader = pypdf.PdfReader(str(pdf_path))
    idx = page_num - 1
    if idx < 0 or idx >= len(reader.pages):
        return None

    # ── attempt 1: largest embedded raster ───────────────────────────────────
    warnings.filterwarnings("ignore", message=".*Lookup Table.*")
    best_blob, best_area = None, 0
    try:
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            page_images = list(reader.pages[idx].images)
        for img_file in page_images:
            try:
                pil = PILImage.open(io.BytesIO(img_file.data))
                area = pil.width * pil.height
                if area > best_area and area > 10_000:
                    best_area, best_blob = area, img_file.data
            except Exception:
                continue
    except Exception:
        pass

    if best_blob is not None:
        try:
            buf = io.BytesIO()
            PILImage.open(io.BytesIO(best_blob)).convert("RGB").save(buf, format="PNG")
            return buf.getvalue()
        except Exception:
            pass

    # ── attempt 2: render full page with pypdfium2 (vector-safe) ─────────────
    try:
        import pypdfium2 as pdfium  # type: ignore
        doc = pdfium.PdfDocument(str(pdf_path))
        page = doc[idx]
        bitmap = page.render(scale=2.0)   # 144 dpi — sharp enough for slides
        pil = bitmap.to_pil()
        buf = io.BytesIO()
        pil.convert("RGB").save(buf, format="PNG")
        return buf.getvalue()
    except Exception:
        pass

    return None


# ── Main slide builder ─────────────────────────────────────────────────────────

def _render_slide(prs, slide_data: dict, idx: int, total: int,
                  pdf_path: Path | None) -> None:
    from pptx.enum.text import PP_ALIGN
    from PIL import Image as PILImage

    blank   = prs.slide_layouts[6]   # truly blank layout
    slide   = prs.slides.add_slide(blank)

    title_text  = slide_data.get("title", f"Slide {idx}")
    tag_text    = (slide_data.get("tag") or "").upper()
    bullets     = slide_data.get("bullets", [])
    narration   = slide_data.get("narration", "")
    image_page  = slide_data.get("image_page")

    is_title_slide = (idx == 1)

    # ── background ───────────────────────────────────────────────────────────
    _set_bg(slide, C_BG)

    if is_title_slide:
        _render_title_slide(slide, title_text, tag_text, bullets,
                             narration, idx, total)
    else:
        _render_content_slide(slide, title_text, tag_text, bullets,
                               narration, idx, total, image_page, pdf_path)


def _render_title_slide(slide, title, tag, bullets, narration, idx, total):
    from pptx.enum.text import PP_ALIGN

    panel_w = int(W * 0.57)   # left dark panel — 57% of slide width
    right_x = panel_w + Inches(0.04)
    right_w = W - right_x

    # ── left panel ────────────────────────────────────────────────────────────
    _add_rect(slide, 0, 0, panel_w, H, C_HEADER)

    # sky-blue top rule (4 px) — subtle depth cue
    _add_rect(slide, 0, 0, panel_w, Inches(0.055), C_SKY)

    # sky-blue bottom rule (4 px)
    _add_rect(slide, 0, H - Inches(0.055), panel_w, Inches(0.055), C_SKY)

    # sky-blue vertical strip at the join
    _add_rect(slide, panel_w, 0, Inches(0.042), H, C_SKY)

    # ── title on left panel ───────────────────────────────────────────────────
    title_top = Inches(1.5)
    title_w   = panel_w - MARGIN_X - Inches(0.5)
    tb = _textbox(slide, MARGIN_X, title_top, title_w, Inches(2.8))
    tb.name = "slide_title"
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    run = p.runs[0]
    run.font.size  = Pt(42 if len(title) < 38 else (34 if len(title) < 55 else 28))
    run.font.bold  = True
    run.font.color.rgb = C_TITLE
    run.font.name  = "Calibri"

    # ── tag / subtitle below title ────────────────────────────────────────────
    tag_y = Inches(4.55)
    if tag:
        tb_tag = _textbox(slide, MARGIN_X, tag_y, title_w, Inches(0.45))
        tb_tag.name = "slide_tag"
        _para(tb_tag.text_frame, tag, 15, bold=False, color=C_SKY)
        tag_y += Inches(0.50)

    # ── key points (up to 4) on left panel ───────────────────────────────────
    if bullets:
        tb3 = _textbox(slide, MARGIN_X, tag_y, title_w, H - tag_y - Inches(0.4))
        tf3 = tb3.text_frame
        tf3.word_wrap = True
        for j, b in enumerate(bullets[:4]):
            _para(tf3, "—  " + b.lstrip(), 14,
                  color=RGBColor(0x86, 0xEF, 0xAC),
                  space_before_pt=(6 if j > 0 else 0))

    # ── right panel — decorative, clean ──────────────────────────────────────
    # single restrained sky-blue accent strip, left of right panel
    _add_rect(slide, right_x + Inches(0.28), Inches(1.6),
              Inches(0.04), Inches(4.3), C_SKY)

    # slide counter — bottom right corner, unobtrusive
    tb4 = _textbox(slide, W - Inches(1.4), H - Inches(0.48),
                   Inches(1.2), Inches(0.36))
    tb4.name = "slide_counter"
    _para(tb4.text_frame, f"{idx} / {total}", 12, color=C_DIM)

    if narration:
        slide.notes_slide.notes_text_frame.text = narration


def _render_content_slide(slide, title, tag, bullets, narration,
                           idx, total, image_page, pdf_path):
    from pptx.enum.text import PP_ALIGN

    # ── header (deep slate, full width) ──────────────────────────────────────
    _add_rect(slide, 0, 0, W, HDR_H, C_HEADER)

    # slight depth band at very top (8 px lighter slate)
    _add_rect(slide, 0, 0, W, Inches(0.08), C_HDR_MID)

    # sky-blue rule at bottom of header
    _add_rect(slide, 0, HDR_H - Inches(0.055), W, Inches(0.055), C_SKY)

    # ── tag (plain sky-blue text, no badge, above title) ─────────────────────
    if tag:
        tb_tag = _textbox(slide, MARGIN_X, MARGIN_Y, W - MARGIN_X * 2 - Inches(1.2), Inches(0.34))
        tb_tag.name = "slide_tag"
        _para(tb_tag.text_frame, tag, 11, bold=True, color=C_SKY)

    # ── title ─────────────────────────────────────────────────────────────────
    title_top = MARGIN_Y + (Inches(0.36) if tag else Inches(0.04))
    title_h   = HDR_H - title_top - Inches(0.10)
    tb_title  = _textbox(slide, MARGIN_X, title_top,
                          W - MARGIN_X - Inches(1.5), title_h)
    tb_title.name = "slide_title"
    tb_title.text_frame.word_wrap = True
    p = tb_title.text_frame.paragraphs[0]
    p.text = title
    font_sz = 32 if len(title) < 42 else (25 if len(title) < 62 else 20)
    run = p.runs[0]
    run.font.size      = Pt(font_sz)
    run.font.bold      = True
    run.font.color.rgb = C_TITLE
    run.font.name      = "Calibri"

    # ── counter — right-aligned in header, unobtrusive ───────────────────────
    tb_cnt = _textbox(slide, W - Inches(1.5), MARGIN_Y + Inches(0.55),
                      Inches(1.3), Inches(0.38))
    tb_cnt.name = "slide_counter"
    _para(tb_cnt.text_frame, f"{idx} / {total}", 13, color=C_DIM)

    # ── sky-blue progress bar at very bottom ──────────────────────────────────
    prog_frac = idx / total
    _add_rect(slide, 0, H - Inches(0.065), W, Inches(0.065), C_RULE)
    _add_rect(slide, 0, H - Inches(0.065), int(W * prog_frac), Inches(0.065), C_SKY)

    # ── content area ──────────────────────────────────────────────────────────
    content_top = HDR_H + Inches(0.28)
    content_bot = H - Inches(0.14)
    content_h   = content_bot - content_top

    img_blob = None
    if image_page and pdf_path:
        img_blob = _extract_pdf_page_image(pdf_path, image_page)

    # slim sky-blue left accent bar
    _add_rect(slide, MARGIN_X - Inches(0.22), content_top,
              ACCENT_W, content_h, C_SKY)

    bullet_x = MARGIN_X + GUTTER
    bullet_w = (IMG_X - MARGIN_X - GUTTER - Inches(0.18)) if img_blob \
               else (W - MARGIN_X - Inches(0.55))

    # ── bullets ───────────────────────────────────────────────────────────────
    if bullets:
        tb_b = _textbox(slide, bullet_x, content_top, bullet_w, content_h)
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True

        n_bullets = len([b for b in bullets if b.strip()])
        base_sz   = 22 if n_bullets <= 4 else (19 if n_bullets <= 6 else 16)

        for j, bullet in enumerate(bullets):
            is_sub = bullet.startswith("  ")
            text   = bullet.lstrip()
            if not text:
                if j > 0:
                    _para(tf_b, "", base_sz - 4, space_before_pt=2)
                continue
            prefix = "     ·  " if is_sub else "—  "
            color  = C_SUB if is_sub else C_TEXT
            sz     = base_sz - 3 if is_sub else base_sz
            _para(tf_b, prefix + text, sz, color=color,
                  space_before_pt=(8 if j > 0 and not is_sub else 2))

    # ── image panel ───────────────────────────────────────────────────────────
    if img_blob:
        try:
            from PIL import Image as PILImage
            pil = PILImage.open(io.BytesIO(img_blob))
            iw, ih = pil.size

            panel_h = int(content_bot - IMG_TOP)
            panel_w = int(IMG_W)
            scale   = min(panel_w / iw, panel_h / ih)
            nw      = max(914, int(iw * scale))
            nh      = max(914, int(ih * scale))

            # sky-blue border → white interior → image
            _add_rect(slide, IMG_X - Inches(0.10), IMG_TOP - Inches(0.10),
                      IMG_W + Inches(0.20), Inches(panel_h / 914400 + 0.20), C_IMGBDR)
            _add_rect(slide, IMG_X - Inches(0.05), IMG_TOP - Inches(0.05),
                      IMG_W + Inches(0.10), Inches(panel_h / 914400 + 0.10), C_IMGBG)

            ox = IMG_X + (IMG_W - nw) // 2
            oy = IMG_TOP + (panel_h - nh) // 2
            slide.shapes.add_picture(io.BytesIO(img_blob), ox, oy, nw, nh)
        except Exception:
            pass

    if narration:
        slide.notes_slide.notes_text_frame.text = narration


def build_pptx(plan: list[dict], out_path: Path,
               pdf_path: Path | None = None) -> None:
    from pptx import Presentation

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    total = len(plan)
    for i, slide_data in enumerate(plan, 1):
        _render_slide(prs, slide_data, i, total, pdf_path)

    prs.save(str(out_path))
    print(f"Saved: {out_path}  ({total} slides)")


